from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def add_card(slide, title, body, x, y, w, h):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(236, 245, 255)
    card.line.color.rgb = RGBColor(78, 123, 205)
    card.line.width = Pt(1.8)

    title_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(y + 0.08), Inches(w - 0.16), Inches(0.42))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(60, 103, 191)
    title_bar.line.fill.background()
    add_textbox(slide, x + 0.14, y + 0.12, w - 0.3, 0.3, title, size=15, bold=True, color=RGBColor(255, 255, 255))

    body_box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.56), Inches(w - 0.28), Inches(h - 0.66))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(25, 32, 55)


def draw_network_house(slide):
    cx = 4.15
    cy = 5.6

    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 1.38), Inches(cy - 1.28), Inches(2.76), Inches(2.76))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(201, 225, 255)
    glow.line.fill.background()

    core = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 1.04), Inches(cy - 0.95), Inches(2.08), Inches(2.08))
    core.fill.solid()
    core.fill.fore_color.rgb = RGBColor(89, 136, 228)
    core.line.color.rgb = RGBColor(255, 255, 255)
    core.line.width = Pt(1.5)

    roof = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - 0.46), Inches(cy - 0.36), Inches(0.92), Inches(0.48))
    roof.fill.solid()
    roof.fill.fore_color.rgb = RGBColor(255, 255, 255)
    roof.line.fill.background()

    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.34), Inches(cy - 0.03), Inches(0.68), Inches(0.5))
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(255, 255, 255)
    body.line.fill.background()

    door = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.05), Inches(cy + 0.2), Inches(0.1), Inches(0.24))
    door.fill.solid()
    door.fill.fore_color.rgb = RGBColor(89, 136, 228)
    door.line.fill.background()

    nodes = [
        (2.0, 4.2, "IOT"),
        (6.28, 4.2, "AI"),
        (2.0, 7.0, "APP"),
        (6.28, 7.0, "MQTT"),
        (4.12, 3.55, "CLOUD"),
        (4.12, 7.7, "DATA"),
    ]

    for nx, ny, label in nodes:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(cy), Inches(nx), Inches(ny)
        )
        connector.line.color.rgb = RGBColor(62, 115, 216)
        connector.line.width = Pt(2.4)

        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - 0.28), Inches(ny - 0.28), Inches(0.56), Inches(0.56))
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(63, 114, 210)
        node.line.color.rgb = RGBColor(255, 255, 255)
        node.line.width = Pt(1.1)
        add_textbox(slide, nx - 0.26, ny - 0.08, 0.52, 0.16, label, size=8.5, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def draw_hero_visual(slide):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.75), Inches(4.02), Inches(2.78), Inches(2.72))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(10, 31, 78)
    panel.fill.transparency = 0.08
    panel.line.color.rgb = RGBColor(119, 233, 255)
    panel.line.width = Pt(2.0)

    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.1), Inches(4.45), Inches(2.1), Inches(1.9))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(69, 215, 255)
    glow.fill.transparency = 0.68
    glow.line.fill.background()

    # Central home hub (sharp vector shapes)
    hub = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(3.63), Inches(5.02), Inches(1.02), Inches(0.86))
    hub.fill.solid()
    hub.fill.fore_color.rgb = RGBColor(55, 142, 255)
    hub.line.color.rgb = RGBColor(208, 247, 255)
    hub.line.width = Pt(1.6)

    roof = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.86), Inches(5.1), Inches(0.56), Inches(0.33))
    roof.fill.solid()
    roof.fill.fore_color.rgb = RGBColor(239, 251, 255)
    roof.line.fill.background()

    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.95), Inches(5.36), Inches(0.38), Inches(0.27))
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(239, 251, 255)
    body.line.fill.background()

    # Ring links + nodes
    cx, cy = 4.14, 5.47
    nodes = [
        (3.08, 4.63, "APP"),
        (5.13, 4.63, "AI"),
        (2.98, 5.95, "CLOUD"),
        (5.21, 5.95, "IOT"),
        (4.14, 6.38, "DATA"),
    ]
    for nx, ny, label in nodes:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(cy), Inches(nx), Inches(ny))
        line.line.color.rgb = RGBColor(98, 221, 255)
        line.line.width = Pt(1.9)

        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - 0.2), Inches(ny - 0.2), Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(42, 117, 228)
        dot.line.color.rgb = RGBColor(214, 247, 255)
        dot.line.width = Pt(1.1)
        add_textbox(slide, nx - 0.35, ny + 0.2, 0.7, 0.16, label, size=7.8, bold=True, color=RGBColor(226, 250, 255), align=PP_ALIGN.CENTER)

    # Extra visual widgets (to avoid plain center)
    widgets = [
        (2.2, 4.35, "MOBILE\nCONTROL"),
        (5.66, 4.35, "SMART\nALERTS"),
        (2.2, 6.18, "ENERGY\nSAVING"),
        (5.66, 6.18, "SAFETY\nMONITOR"),
    ]
    for x, y, txt in widgets:
        w = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.08), Inches(0.72))
        w.fill.solid()
        w.fill.fore_color.rgb = RGBColor(23, 66, 148)
        w.fill.transparency = 0.06
        w.line.color.rgb = RGBColor(104, 230, 255)
        w.line.width = Pt(1.2)
        add_textbox(slide, x + 0.04, y + 0.11, 1.0, 0.5, txt, size=7.5, bold=True, color=RGBColor(227, 250, 255), align=PP_ALIGN.CENTER)


def draw_house_network_with_four_boxes(slide):
    cx, cy = 4.14, 6.85

    center_glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 1.25), Inches(cy - 1.15), Inches(2.5), Inches(2.3))
    center_glow.fill.solid()
    center_glow.fill.fore_color.rgb = RGBColor(79, 226, 255)
    center_glow.fill.transparency = 0.72
    center_glow.line.fill.background()

    core = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.8), Inches(cy - 0.75), Inches(1.6), Inches(1.5))
    core.fill.solid()
    core.fill.fore_color.rgb = RGBColor(44, 124, 236)
    core.line.color.rgb = RGBColor(210, 249, 255)
    core.line.width = Pt(1.8)

    roof = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - 0.35), Inches(cy - 0.34), Inches(0.7), Inches(0.36))
    roof.fill.solid()
    roof.fill.fore_color.rgb = RGBColor(245, 253, 255)
    roof.line.fill.background()

    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.24), Inches(cy - 0.02), Inches(0.48), Inches(0.35))
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(245, 253, 255)
    body.line.fill.background()

    door = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.05), Inches(cy + 0.12), Inches(0.1), Inches(0.18))
    door.fill.solid()
    door.fill.fore_color.rgb = RGBColor(44, 124, 236)
    door.line.fill.background()

    box_data = [
        (0.65, 5.65, 2.2, 1.2, "AKILLI KONTROL", "Mobil uygulama ile cihazlari\nanlik yonetme ve izleme"),
        (5.42, 5.65, 2.2, 1.2, "ONERI MOTORU", "Kullanici aliskanliklarina gore\nproaktif eylem onerileri"),
        (0.65, 8.05, 2.2, 1.2, "ENERJI VERIMLILIGI", "Gereksiz tuketimi azaltan\nsenaryo ve otomasyonlar"),
        (5.42, 8.05, 2.2, 1.2, "GUVENLIK", "Anomali tespiti ve kritik\nbildirimler ile guvenli yasam"),
    ]

    anchors = [
        (2.85, 6.25),
        (5.42, 6.25),
        (2.85, 8.25),
        (5.42, 8.25),
    ]

    for i, (x, y, w, h, title, body_text) in enumerate(box_data):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(221, 248, 255)
        card.fill.transparency = 0.04
        card.line.color.rgb = RGBColor(118, 237, 255)
        card.line.width = Pt(1.6)

        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(y + 0.08), Inches(w - 0.16), Inches(0.34))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(80, 212, 245)
        band.line.fill.background()
        add_textbox(slide, x + 0.12, y + 0.12, w - 0.22, 0.22, title, size=11, bold=True, color=RGBColor(8, 33, 78), align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.12, y + 0.47, w - 0.22, h - 0.52, body_text, size=10, color=RGBColor(21, 34, 49), align=PP_ALIGN.CENTER)

        ax, ay = anchors[i]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(cy), Inches(ax), Inches(ay))
        conn.line.color.rgb = RGBColor(96, 228, 255)
        conn.line.width = Pt(2.4)


def add_dark_glow_background(slide):
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(8.27), Inches(11.69))
    base.fill.solid()
    base.fill.fore_color.rgb = RGBColor(8, 25, 68)
    base.line.fill.background()

    right_tint = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.9), Inches(0), Inches(4.37), Inches(11.69))
    right_tint.fill.solid()
    right_tint.fill.fore_color.rgb = RGBColor(18, 132, 171)
    right_tint.fill.transparency = 0.32
    right_tint.line.fill.background()

    top_glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(-0.9), Inches(6.3), Inches(3.2))
    top_glow.fill.solid()
    top_glow.fill.fore_color.rgb = RGBColor(65, 201, 239)
    top_glow.fill.transparency = 0.72
    top_glow.line.fill.background()

    center_glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.1), Inches(3.55), Inches(4.0), Inches(3.0))
    center_glow.fill.solid()
    center_glow.fill.fore_color.rgb = RGBColor(92, 221, 255)
    center_glow.fill.transparency = 0.78
    center_glow.line.fill.background()


def add_dark_card(slide, title, body, x, y, w, h):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(224, 250, 255)
    card.fill.transparency = 0.06
    card.line.color.rgb = RGBColor(155, 245, 255)
    card.line.width = Pt(1.6)

    title_band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.42))
    title_band.fill.solid()
    title_band.fill.fore_color.rgb = RGBColor(95, 227, 243)
    title_band.line.fill.background()

    add_textbox(slide, x + 0.15, y + 0.13, w - 0.35, 0.3, title, size=15, bold=True, color=RGBColor(7, 27, 58))

    body_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.56), Inches(w - 0.3), Inches(h - 0.67))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(18, 27, 40)


def main():
    prs = Presentation()
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_dark_glow_background(slide)

    top = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(0), Inches(1.3), Inches(8.27), Inches(1.25))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(28, 108, 199)
    top.fill.transparency = 0.15
    top.line.fill.background()

    low = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(0), Inches(9.72), Inches(8.27), Inches(1.15))
    low.fill.solid()
    low.fill.fore_color.rgb = RGBColor(18, 165, 202)
    low.fill.transparency = 0.1
    low.line.fill.background()

    add_textbox(slide, 0.35, 0.14, 5.2, 0.35, "DOKUZ EYLUL UNIVERSITY", size=24, bold=True, color=RGBColor(238, 247, 255))
    add_textbox(slide, 6.52, 0.16, 1.4, 0.28, "#20250524", size=14, bold=True, color=RGBColor(238, 247, 255), align=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.35, 0.5, 3.8, 0.55, "SYNAPSE", size=44, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, 4.45, 0.57, 3.45, 0.5, "Durukan Hacioglu, Berkay Karakus\nMentor: Prof. Dr. Yunus Dogan", size=11, color=RGBColor(229, 244, 255), align=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.35, 1.02, 7.6, 0.35, "PERSONALIZED SMART HOME INTERACTION WITH IOT DEVICES", size=17, bold=True, color=RGBColor(255, 255, 255))

    intro_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(2.2), Inches(7.57), Inches(1.5))
    intro_card.fill.solid()
    intro_card.fill.fore_color.rgb = RGBColor(221, 248, 255)
    intro_card.fill.transparency = 0.05
    intro_card.line.color.rgb = RGBColor(120, 238, 255)
    intro_card.line.width = Pt(1.6)
    add_textbox(slide, 0.55, 2.35, 7.2, 0.32, "INTRODUCTION", size=18, bold=True, color=RGBColor(8, 33, 78), align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        0.55,
        2.72,
        7.2,
        0.9,
        "SYNAPSE, IoT tabanli akilli yasam ortamlarinda kullanici davranislarini analiz ederek\nproaktif ve kisisellestirilmis oneriler sunan bir sistemdir. Mobil uygulama, FastAPI servisleri\nve gercek zamanli bildirim altyapisi ile konforu artirirken enerji verimliligini ve guvenligi\ngelistirir. Sistem, kullanici onayi ile calisan aciklanabilir bir karar mekanizmasi sunar.",
        size=11.5,
        color=RGBColor(20, 34, 52),
        align=PP_ALIGN.CENTER,
    )

    draw_house_network_with_four_boxes(slide)

    footer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(10.35), Inches(7.35), Inches(0.85))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(13, 52, 125)
    footer.fill.transparency = 0.04
    footer.line.color.rgb = RGBColor(80, 216, 245)
    footer.line.width = Pt(1.8)
    add_textbox(
        slide,
        0.65,
        10.47,
        6.95,
        0.45,
        "Genc Beyinler Yeni Fikirler\nProje Pazari ve Bitirme Projeleri Ortak Sergisi",
        size=14.5,
        bold=True,
        color=RGBColor(236, 249, 255),
        align=PP_ALIGN.CENTER,
    )

    prs.save("Synapse_Poster_A4_v6.pptx")


if __name__ == "__main__":
    main()
